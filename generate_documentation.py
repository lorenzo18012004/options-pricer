"""
Génère la documentation complète du projet Pricer :
- Pricer_Documentation_Complete.docx : document Word détaillé
- Pricer_Resume.pptx : PowerPoint de synthèse
"""

from pathlib import Path

OUTPUT_DIR = Path(__file__).parent / "docs"
OUTPUT_DIR.mkdir(exist_ok=True)


def create_word_document():
    """Crée le document Word complet."""
    from docx import Document
    from docx.shared import Pt, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    doc = Document()
    
    # Style titre principal
    title = doc.add_heading("Options Pricer - Documentation Complète", 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_paragraph(
        "Projet développé from scratch. Ce document explique l'intégralité des choix "
        "techniques, des modèles et des conventions utilisés."
    )
    
    # ========== 1. INTRODUCTION ==========
    doc.add_heading("1. Introduction et objectifs", level=1)
    doc.add_paragraph(
        "Ce pricer d'options a été développé dans le cadre d'un projet personnel visant à "
        "maîtriser les modèles de pricing d'options et leur implémentation. L'objectif est "
        "de fournir un outil pédagogique et démonstratif couvrant plusieurs classes de "
        "produits : options vanilles, stratégies de volatilité, options barrières et swaps de taux."
    )
    doc.add_paragraph(
        "Le projet repose sur des données de marché live (Yahoo Finance) et implémente "
        "les modèles de référence utilisés en finance quantitative."
    )
    
    # ========== 2. ARCHITECTURE ==========
    doc.add_heading("2. Architecture du projet", level=1)
    doc.add_paragraph("L'application est structurée en modules distincts :")
    
    arch_items = [
        ("core/", "Cœur quantitatif : Black-Scholes, Heston, solveur IV, courbes de taux"),
        ("models/", "Modèles numériques : arbre binomial, Monte Carlo, surfaces SVI"),
        ("instruments/", "Produits : options vanilles, barrières, swaps"),
        ("data/", "Connecteur Yahoo Finance, nettoyage des chaînes d'options"),
        ("services/", "Services métier : chargement marché, construction expirations"),
        ("pages/", "Interface Streamlit : une page par produit"),
    ]
    for name, desc in arch_items:
        p = doc.add_paragraph(style='List Bullet')
        p.add_run(name).bold = True
        p.add_run(f" — {desc}")
    
    # ========== 3. DONNÉES ==========
    doc.add_heading("3. Sources de données et conventions", level=1)
    
    doc.add_heading("3.1 Yahoo Finance", level=2)
    doc.add_paragraph(
        "Yahoo Finance (via yfinance) a été choisi comme unique source de données pour plusieurs raisons :"
    )
    doc.add_paragraph(
        "• Gratuit et sans clé API : idéal pour un projet personnel et démo",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• Données US complètes : spot, chaînes d'options (bid/ask, IV, volume, OI), historiques",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• Courbe des taux : 4 points Treasury (^IRX 3M, ^FVX 5Y, ^TNX 10Y, ^TYX 30Y)",
        style='List Bullet'
    )
    doc.add_paragraph(
        "Limitation : en production, un desk utiliserait Bloomberg ou Refinitiv pour des données "
        "institutionnelles et une meilleure qualité."
    )
    
    doc.add_heading("3.2 Conventions de marché", level=2)
    doc.add_paragraph(
        "Pour les options equity vanilles, les conventions desk ont été adoptées :"
    )
    doc.add_paragraph(
        "• Horloge de maturité : jours calendaires (calendar days), pas jours ouvrés",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• Day count : ACT/365 pour l'actualisation et le calcul de T",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• T = jours_restants / 365",
        style='List Bullet'
    )
    doc.add_paragraph(
        "Ces conventions sont standard pour les options equity sur actions US."
    )
    
    doc.add_heading("3.3 Fenêtre de volatilité historique", level=2)
    doc.add_paragraph(
        "La volatilité historique est calculée sur une fenêtre adaptée à la maturité de l'option :"
    )
    doc.add_paragraph(
        "• < 14 jours ouvrés → 10 jours ; < 45j → 20j ; < 90j → 60j ; < 180j → 120j ; sinon 252j",
        style='List Bullet'
    )
    doc.add_paragraph(
        "Cela permet d'aligner la volatilité réalisée avec l'horizon de l'option."
    )
    
    # ========== 4. OPTIONS VANILLES ==========
    doc.add_heading("4. Options vanilles - Modèles de pricing", level=1)
    
    doc.add_heading("4.1 Black-Scholes-Merton (BSM)", level=2)
    doc.add_paragraph(
        "Modèle de référence pour les options européennes. Choix justifié par :"
    )
    doc.add_paragraph(
        "• Formules fermées : pricing et Greeks analytiques, rapides et stables",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• Support des dividendes continus (dividend yield q) via Merton (1973)",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• Base pour le solveur de volatilité implicite (IV) et le benchmark Monte Carlo",
        style='List Bullet'
    )
    doc.add_paragraph(
        "Formules : d1 = [ln(S/K) + (r - q + 0.5σ²)T] / (σ√T), d2 = d1 - σ√T. "
        "Call = S·e^(-qT)·N(d1) - K·e^(-rT)·N(d2)."
    )
    
    doc.add_heading("4.2 Greeks (1er et 2nd ordre)", level=2)
    doc.add_paragraph(
        "Les Greeks sont implémentés pour le hedging et l'analyse de risque :"
    )
    doc.add_paragraph(
        "• 1er ordre : Delta, Gamma, Vega, Theta, Rho — sensibilités de base",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• 2nd ordre : Vanna (dDelta/dVol), Volga (dVega/dVol), Charm (dDelta/dT), Speed, Color, Zomma",
        style='List Bullet'
    )
    doc.add_paragraph(
        "Vanna et Charm sont essentiels pour le P&L attribution (décomposition du P&L en composantes)."
    )
    
    doc.add_heading("4.3 Solveur de volatilité implicite (IV)", level=2)
    doc.add_paragraph(
        "Pour retrouver l'IV à partir d'un prix de marché :"
    )
    doc.add_paragraph(
        "• Newton-Raphson en méthode principale (convergence quadratique, utilise la Vega comme dérivée)",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• Fallback Brent : robuste, garantit la convergence si une solution existe",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• Vérification des bornes no-arbitrage avant résolution",
        style='List Bullet'
    )
    
    doc.add_heading("4.4 Modèle de Heston (1993)", level=2)
    doc.add_paragraph(
        "Modèle à volatilité stochastique. Choix justifié par :"
    )
    doc.add_paragraph(
        "• Génère naturellement le smile et le skew de volatilité (contrairement à BSM à vol plate)",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• Capture l'effet de levier : corrélation spot-vol négative (rho < 0) pour les actions",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• Mean-reversion de la variance : kappa (vitesse), theta (niveau long terme)",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• Solution semi-analytique via inversion de Fourier (Lewis 2000) — rapide vs Monte Carlo",
        style='List Bullet'
    )
    doc.add_paragraph(
        "Processus : dS = (r-q)S dt + √v S dW₁ ; dv = κ(θ-v) dt + ξ√v dW₂ ; Corr(dW₁,dW₂) = ρ"
    )
    
    doc.add_heading("4.5 Monte Carlo", level=2)
    doc.add_paragraph(
        "Simulations de trajectoires pour valider les modèles analytiques et pricer les path-dependants :"
    )
    doc.add_paragraph(
        "• Schéma d'Euler pour dS = (r-q-½σ²)S dt + σS dW",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• Antithetic variates : réduction de variance en utilisant -Z pour chaque Z",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• Seed fixe pour reproductibilité",
        style='List Bullet'
    )
    
    doc.add_heading("4.6 Arbre binomial (Cox-Ross-Rubinstein)", level=2)
    doc.add_paragraph(
        "Pour les options américaines (exercice anticipé) :"
    )
    doc.add_paragraph(
        "• CRR : u = exp(σ√Δt), d = 1/u, probabilité risque-neutre p = (exp((r-q)Δt) - d)/(u - d)",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• À chaque nœud : max(valeur continuation, valeur d'exercice immédiat)",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• 200 steps par défaut pour le pricing américain (EU vs US)",
        style='List Bullet'
    )
    
    doc.add_heading("4.7 SVI (Stochastic Volatility Inspired)", level=2)
    doc.add_paragraph(
        "Paramétrisation du smile de volatilité pour lisser les IV de marché :"
    )
    doc.add_paragraph(
        "• σ²(k) = a + b·(ρ(k-m) + √((k-m)² + σ²)) avec k = log(K/F)",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• Calibration par curve_fit (scipy) avec contraintes d'arbitrage",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• Option « Use SVI smoothed IV » : utilise l'IV SVI au lieu de l'IV brute pour pricing et Greeks",
        style='List Bullet'
    )
    
    # ========== 5. STRATÉGIES VOLATILITÉ ==========
    doc.add_heading("5. Stratégies de volatilité (Straddle, Strangle)", level=1)
    doc.add_paragraph(
        "Straddle = Call ATM + Put ATM. Strangle = Call OTM + Put OTM (même delta). "
        "Stratégies vega-long : profitent d'une hausse de la volatilité."
    )
    doc.add_paragraph(
        "Pricing via BSM avec une vol unique (IV du call ou moyenne). Analyse P&L attribution, "
        "convergence binomiale, comparaison Heston, heatmap Spot/Vol, simulation de delta-hedging."
    )
    
    # ========== 6. OPTIONS BARRIÈRES ==========
    doc.add_heading("6. Options barrières", level=1)
    doc.add_paragraph(
        "Options dont le payoff dépend du franchissement d'un niveau (barrière) :"
    )
    doc.add_paragraph(
        "• Knock-out : l'option disparaît si la barrière est touchée",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• Knock-in : l'option s'active si la barrière est touchée",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• Down/Up : direction de la barrière par rapport au spot",
        style='List Bullet'
    )
    doc.add_paragraph(
        "Pricing par Monte Carlo avec ajustement Brownian bridge : pour détecter un franchissement "
        "entre deux pas de temps discrets, on utilise la propriété du pont brownien pour estimer "
        "la probabilité de touch. Plus précis qu'un check naïf aux seuls points discrets."
    )
    
    # ========== 7. SWAPS DE TAUX ==========
    doc.add_heading("7. Swaps de taux d'intérêt", level=1)
    doc.add_paragraph(
        "Swap vanille : échange de flux fixe contre flux flottant (Libor/SOFR)."
    )
    doc.add_heading("7.1 Courbe des taux", level=2)
    doc.add_paragraph(
        "• Source : 4 points Treasury Yahoo (3M, 5Y, 10Y, 30Y) — par yields",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• Bootstrap : construction de la zero-curve à partir des par yields (obligations au pair)",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• Hypothèse : coupons semi-annuels (payment_freq=2), prix pair = 1",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• Interpolation : log-linéaire des discount factors entre les nœuds",
        style='List Bullet'
    )
    
    doc.add_heading("7.2 Pricing du swap", level=2)
    doc.add_paragraph(
        "• Jambe fixe : PV = Σ (Notional × FixedRate × αᵢ × DF(tᵢ))",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• Jambe flottante (vanille) : PV = Notional × (1 - DF(T))",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• NPV (payer) = PV(flottant) - PV(fixe)",
        style='List Bullet'
    )
    
    doc.add_heading("7.3 Risque et hedge", level=2)
    doc.add_paragraph(
        "• DV01 : sensibilité à un mouvement parallèle de 1 bp de la courbe",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• Par rate : taux fixe qui annule le NPV",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• Key-rate DV01 : sensibilité par point de la courbe",
        style='List Bullet'
    )
    
    # ========== 8. QUALITÉ ET TESTS ==========
    doc.add_heading("8. Qualité des données et tests", level=1)
    doc.add_heading("8.1 Contrôles no-arbitrage", level=2)
    doc.add_paragraph(
        "• Monotonie des prix (call croissant en K, put décroissant)",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• Convexité (butterfly) : pas de convexité négative",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• Put-call parity : C - P = S·e^(-qT) - K·e^(-rT)",
        style='List Bullet'
    )
    doc.add_paragraph(
        "• Filtrage par moneyness (0.80-1.20) pour se concentrer sur les options liquides",
        style='List Bullet'
    )
    
    doc.add_heading("8.2 Tests unitaires (42 tests)", level=2)
    doc.add_paragraph(
        "Couverture : Black-Scholes (parité, bornes, Greeks), IV solver, binomial, Monte Carlo, "
        "Heston, barrières, courbe de taux, swap, DataCleaner, SVI. Exécution : pytest."
    )
    
    # ========== 9. STACK TECHNIQUE ==========
    doc.add_heading("9. Stack technique", level=1)
    doc.add_paragraph(
        "Python 3, NumPy, SciPy, Pandas, yfinance, Streamlit, Plotly. Tests : pytest."
    )
    
    # ========== 10. CONCLUSION ==========
    doc.add_heading("10. Conclusion", level=1)
    doc.add_paragraph(
        "Ce pricer démontre une maîtrise des modèles fondamentaux de la finance quantitative "
        "et de leur implémentation. Il couvre les produits principaux (vanilles, barrières, swaps) "
        "avec des conventions de marché cohérentes et des contrôles de qualité. "
        "Projet pédagogique et portfolio, conçu pour illustrer les concepts de pricing d'options."
    )
    
    out_path = OUTPUT_DIR / "Pricer_Documentation_Complete.docx"
    doc.save(out_path)
    print(f"Word créé : {out_path}")
    return out_path


def create_powerpoint():
    """Crée le PowerPoint de synthèse — design type finance/académique (fond sombre, sobre)."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor

    # Couleurs : fond navy sombre, texte clair, accent doré (style "Crise des Modèles")
    BG_DARK = RGBColor(0x0F, 0x17, 0x2A)   # Navy très sombre
    TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    TEXT_LIGHT = RGBColor(0xE2, 0xE8, 0xF0)  # Gris clair
    ACCENT = RGBColor(0xCA, 0x8A, 0x04)      # Or / ambre discret

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    def set_dark_background(slide):
        """Fond sombre sur la slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK

    def add_title_slide(title, subtitle="", slide_num=None, total=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_dark_background(slide)
        # Numéro de slide (style référence)
        if slide_num and total:
            num_tx = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.4))
            num_tx.text_frame.paragraphs[0].text = f"— {slide_num} of {total} —"
            num_tx.text_frame.paragraphs[0].font.size = Pt(11)
            num_tx.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
            num_tx.text_frame.paragraphs[0].alignment = 1
        # Ligne d'accent horizontale en haut
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(2), Emu(500))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT
        line.line.fill.background()
        # Titre centré
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1.8))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title.upper()
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.alignment = 1  # Center
        if subtitle:
            tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(1.2))
            tf2 = tx2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.font.size = Pt(18)
            p2.font.color.rgb = TEXT_LIGHT
            p2.alignment = 1
        return slide

    def add_content_slide(title, bullets, slide_num=None, total=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_dark_background(slide)
        if slide_num and total:
            num_tx = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.4))
            num_tx.text_frame.paragraphs[0].text = f"— {slide_num} of {total} —"
            num_tx.text_frame.paragraphs[0].font.size = Pt(11)
            num_tx.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
            num_tx.text_frame.paragraphs[0].alignment = 1
        # Ligne d'accent sous le titre
        tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(8.8), Inches(0.9))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        # Ligne horizontale accent
        line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.35), Inches(1.5), Emu(400))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT
        line.line.fill.background()
        # Contenu
        tx2 = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(8.8), Inches(5.3))
        tf2 = tx2.text_frame
        tf2.word_wrap = True
        for i, b in enumerate(bullets):
            para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            para.text = "•  " + b
            para.font.size = Pt(15)
            para.font.color.rgb = TEXT_LIGHT
            para.space_before = Pt(12)
            para.level = 0
        return slide

    TOTAL = 11
    # Slide 1
    add_title_slide("Options Pricer", "Documentation complète — Résumé exécutif", 1, TOTAL)
    # Slide 2
    add_content_slide("1. Objectif du projet", [
        "Pricer d'options développé from scratch",
        "Couvre : vanilles, stratégies vol, barrières, swaps",
        "Données live Yahoo Finance",
        "Projet pédagogique et portfolio"
    ], 2, TOTAL)
    # Slide 3
    add_content_slide("2. Architecture", [
        "core/ : Black-Scholes, Heston, IV solver, courbes",
        "models/ : Binomial, Monte Carlo, SVI",
        "instruments/ : Options, barrières, swaps",
        "data/ : Yahoo Finance, nettoyage",
        "pages/ : Interface Streamlit"
    ], 3, TOTAL)
    # Slide 4
    add_content_slide("3. Conventions de marché", [
        "Maturité : jours calendaires",
        "Day count : ACT/365",
        "T = jours_restants / 365",
        "Vol historique : fenêtre adaptée à la maturité"
    ], 4, TOTAL)
    # Slide 5
    add_content_slide("4. Modèles de pricing — Options vanilles", [
        "Black-Scholes-Merton : européennes, dividendes (q)",
        "Heston : volatilité stochastique, smile/skew",
        "Monte Carlo : Euler + antithetic variates",
        "Binomial CRR : américaines (200 steps)",
        "SVI : lissage du smile de volatilité"
    ], 5, TOTAL)
    # Slide 6
    add_content_slide("5. Greeks et IV", [
        "1er ordre : Delta, Gamma, Vega, Theta, Rho",
        "2nd ordre : Vanna, Volga, Charm, Speed, Color, Zomma",
        "Solveur IV : Newton-Raphson + Brent",
        "P&L attribution : décomposition Delta, Gamma, Vega, Theta, Vanna"
    ], 6, TOTAL)
    # Slide 7
    add_content_slide("6. Options barrières", [
        "Knock-in / Knock-out, Down / Up",
        "Pricing : Monte Carlo + Brownian bridge",
        "Détection du franchissement intra-step",
        "Parité in + out = vanilla"
    ], 7, TOTAL)
    # Slide 8
    add_content_slide("7. Swaps de taux", [
        "Courbe : 4 points Treasury (3M, 5Y, 10Y, 30Y)",
        "Bootstrap par yields → zero-curve",
        "NPV, DV01, Par rate, Key-rate risk",
        "Hedge sizing"
    ], 8, TOTAL)
    # Slide 9
    add_content_slide("8. Qualité et tests", [
        "Contrôles no-arbitrage : monotonie, convexité, put-call parity",
        "42 tests pytest : BSM, IV, binomial, MC, Heston, barrières, courbe, swap, DataCleaner, SVI",
        "Filtrage moneyness 0.80-1.20"
    ], 9, TOTAL)
    # Slide 10
    add_content_slide("9. Stack technique", [
        "Python, NumPy, SciPy, Pandas",
        "yfinance (données)",
        "Streamlit (UI), Plotly (graphiques)",
        "pytest (tests)"
    ], 10, TOTAL)
    # Slide 11
    add_title_slide("Merci", "Documentation complète disponible dans le Word associé", 11, TOTAL)
    
    out_path = OUTPUT_DIR / "Pricer_Resume.pptx"
    prs.save(out_path)
    print(f"PowerPoint créé : {out_path}")
    return out_path


def main():
    print("Génération de la documentation...")
    create_word_document()
    create_powerpoint()
    print(f"\nFichiers générés dans : {OUTPUT_DIR.absolute()}")


if __name__ == "__main__":
    main()
